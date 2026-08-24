# -*- coding: utf-8 -*-
"""统一文档解析：word/xlsx/pptx → 标准化 ParsedDoc。

设计要点：
- 解析是确定性的，不依赖 LLM、不依赖需求；标准文档与学生文件用同一解析器，
  保证比较基线一致（同一解析器 → 同构比较）。
- ParsedDoc 承载 §4.6 scope 定位所需的元素清单：paragraphs/tables/images/
  comments/formulas/header/footer/sections，供 locate.py 与 primitives.py 使用。
- 解析异常收敛为 ParseError；单点提取失败不影响整体（best-effort）。

用法：
    from parsers import parse_file, parse_bytes
    doc = parse_file("标准文档.docx")        # 按扩展名推断类型
    doc = parse_bytes(data, "word")
"""
from __future__ import annotations

import zipfile
from io import BytesIO
from pathlib import Path

try:
    from lxml import etree
except Exception:  # pragma: no cover
    etree = None

ZIP_MAX_ENTRIES = 200
ZIP_MAX_UNCOMPRESSED = 50 * 1024 * 1024


class ParseError(Exception):
    pass


def _cm(emu):
    """EMU（1/914400 英寸）→ 厘米，保留 2 位；无法转换返回 None。"""
    try:
        return round(float(emu) / 360000.0, 2)
    except (TypeError, ValueError):
        return None


def _hl_color(run):
    """run 的突出显示颜色（w:highlight val），无则 None。"""
    try:
        rPr = run._element.rPr
        if rPr is None:
            return None
        h = rPr.find("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}highlight")
        return h.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val") if h is not None else None
    except Exception:
        return None


def _run_color(run):
    try:
        c = run.font.color
        if c is None or c.type is None:
            return None
        return str(c.rgb) if c.rgb is not None else (c.theme_color if hasattr(c, "theme_color") else str(c.type))
    except Exception:
        return None


def _run_font_names(run):
    """run 的字体名集合：显式 w:rFonts 的 ascii/hAnsi/eastAsia 与主题字体名。

    修正（2026-08-19）：原实现只取 run.font.name（ascii 字体），中文字体常写在
    w:rFonts @w:eastAsia（宋体/黑体等）或 @w:asciiTheme/@w:hAnsiTheme（主题字体
    如 minorEastAsia，默认映射宋体）。漏读导致"正文宋体"这类需求无法确定性判定。
    """
    names = set()
    try:
        rPr = run._element.rPr
        if rPr is None:
            return names
        q = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
        rf = rPr.find(q + "rFonts")
        if rf is None:
            return names
        for attr in ("ascii", "hAnsi", "eastAsia",
                     "asciiTheme", "hAnsiTheme", "eastAsiaTheme", "cstheme"):
            v = rf.get(q + attr)
            if v:
                names.add(v)
        if run.font.name:
            names.add(run.font.name)
    except Exception:
        pass
    return names


def _para_border_vals(p):
    """段落边框：返回各边 val 的 dict 或 None。"""
    try:
        pPr = p._p.pPr
        if pPr is None:
            return None
        b = pPr.find("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}pBdr")
        if b is None:
            return None
        out = {}
        for side in ("top", "bottom", "left", "right"):
            el = b.find("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}%s" % side)
            if el is not None:
                out[side] = el.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val")
        return out or None
    except Exception:
        return None


def _para_shading(p):
    """段落底纹：{fill, color, themeColor, themeTint, themeFill, themeFillTint} 或 None。

    修正（2026-08-19）：补 themeFill/themeFillTint —— 需求常写"红色，强调文字颜色2，淡色80%"
    这类主题色底纹，落在 themeFill/themeFillTint 而非 fill，原实现漏读导致只看到 fill 缺失。
    """
    try:
        pPr = p._p.pPr
        if pPr is None:
            return None
        s = pPr.find("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}shd")
        if s is None:
            return None
        q = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
        keys = ("fill", "color", "themeColor", "themeTint",
                "themeFill", "themeFillTint")
        return {k: s.get(q + k) for k in keys if s.get(q + k) is not None} or None
    except Exception:
        return None


def _para_flag(p, tag):
    """段落 pPr 是否含指定布尔控制标签（w:widowControl / w:topLinePunct 等）。

    返回 bool；解析失败返回 None（区别于明确的 False，避免与"显式关闭"混淆时硬判）。
    注意：w:widowControl 无 val 时表示开启(true)，val="0" 表示关闭(false)。
    """
    try:
        pPr = p._p.pPr
        if pPr is None:
            return None
        q = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
        el = pPr.find(q + tag)
        if el is None:
            return False
        val = el.get(q + "val")
        # 显式 val="0" 表示关闭；无 val / val="1"（或其它真值）视为开启
        return not (val == "0")
    except Exception:
        return None


def _align_name(align):
    if align is None:
        return None
    try:
        return {
            0: "left", 1: "center", 2: "right", 3: "justify",
        }.get(int(align), str(align))
    except (TypeError, ValueError):
        return None


def _check_zip(data, label):
    """zip 容器防护：非 zip、条目数或解压体积超限抛 ParseError。"""
    try:
        zf = zipfile.ZipFile(BytesIO(data))
    except (zipfile.BadZipFile, OSError) as e:
        raise ParseError("%s 不是有效的 zip 容器: %s" % (label, e)) from e
    infos = zf.infolist()
    if len(infos) > ZIP_MAX_ENTRIES:
        raise ParseError("%s 条目数 %d 超过上限 %d" % (label, len(infos), ZIP_MAX_ENTRIES))
    declared = sum(i.file_size for i in infos)
    if declared > ZIP_MAX_UNCOMPRESSED:
        raise ParseError("%s 声明解压大小超限" % label)
    real = 0
    for info in infos:
        try:
            with zf.open(info) as f:
                while True:
                    chunk = f.read(1 << 20)
                    if not chunk:
                        break
                    real += len(chunk)
                    if real > ZIP_MAX_UNCOMPRESSED:
                        raise ParseError("%s 实际解压体积超限" % label)
        except ParseError:
            raise
        except (zipfile.BadZipFile, OSError, RuntimeError):
            pass  # 单条目读取失败不中断
    return zf


class ParsedDoc:
    """标准化解析结果。

    字段（全部可空/可缺省，best-effort）：
    - file_type: word | xlsx | pptx
    - text: 全文纯文本
    - paragraphs: [{index, text, style, align, font_name, font_size, bold,
                    italic, underline, first_line_indent, left_indent,
                    space_before, space_after, line_spacing}]
    - tables: [{index, rows, cols, align, header_row}]
    - images: [{index, width_cm, height_cm}]
    - comments: [文本]
    - formulas: 数量（word 内联/块公式个数）
    - footnotes / endnotes: 数量
    - header_text / footer_text: 首页节页眉/页脚文本
    - sections: [{margin_top/bottom/left/right_cm, paper_w/h_cm,
                  orientation, gutter_cm, page_number_fmt, header/footer}]
    - page_breaks: 分页符数量
    - track_changes: bool（word 修订）
    - watermark: bool（word 水印）
    - toc: bool（word 目录域）
    - meta: 原始库对象引用（doc/openpyxl workbook/pptx presentation），
            供 LLM 判定路径按需深入读取
    """

    def __init__(self, file_type, paragraphs=None, tables=None, images=None,
                 comments=None, formulas=0, footnotes=0, endnotes=0,
                 header_text="", footer_text="", sections=None, page_breaks=0,
                 track_changes=False, watermark=False, toc=False,
                 sheet_count=0, slide_count=0, footer_page_field=False,
                 meta=None):
        self.file_type = file_type
        self.paragraphs = paragraphs or []
        self.tables = tables or []
        self.images = images or []
        self.comments = comments or []
        self.formulas = formulas
        self.footnotes = footnotes
        self.endnotes = endnotes
        self.header_text = header_text or ""
        self.footer_text = footer_text or ""
        self.sections = sections or []
        self.page_breaks = page_breaks
        self.track_changes = track_changes
        self.watermark = watermark
        self.toc = toc
        self.sheet_count = sheet_count
        self.slide_count = slide_count
        self.footer_page_field = footer_page_field
        self.meta = meta

    @property
    def text(self):
        return "\n".join(p.get("text", "") for p in self.paragraphs)


# --------------------------------------------------------------------------
# word
# --------------------------------------------------------------------------

def _xml_root(data):
    if etree is None:
        return None
    try:
        return etree.fromstring(data)
    except Exception:
        return None


def _count_xpath(root, path, ns):
    if root is None:
        return 0
    try:
        return len(root.findall(path, ns))
    except Exception:
        return 0


_W_NS = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
         "ns0": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
         "m": "http://schemas.openxmlformats.org/officeDocument/2006/math"}


def _word_part(zf, name):
    try:
        return zf.read(name)
    except KeyError:
        return None


def _parse_word(zf, data):
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document(BytesIO(data))
    paragraphs = []
    for i, p in enumerate(doc.paragraphs):
        run = next((r for r in p.runs if r.text and r.text.strip()), None)
        para = {
            "index": i,
            "text": p.text,
            "style": p.style.name if p.style else None,
            "align": _align_name(p.alignment),
        }
        if run is not None:
            para["font_name"] = _run_font_names(run) or None
            para["font_size"] = (run.font.size.pt if run.font.size else None)
            para["bold"] = bool(run.bold)
            para["italic"] = bool(run.italic)
            para["underline"] = bool(run.underline)
            para["font_color"] = _run_color(run)
            para["highlight"] = _hl_color(run)
        try:
            pf = p.paragraph_format
            para["first_line_indent"] = _cm(pf.first_line_indent)
            para["left_indent"] = _cm(pf.left_indent)
            para["space_before"] = (pf.space_before.pt if pf.space_before else None)
            para["space_after"] = (pf.space_after.pt if pf.space_after else None)
            ls = pf.line_spacing
            para["line_spacing"] = (ls.pt if hasattr(ls, "pt") else ls) if ls else None
        except Exception:
            pass
        para["border"] = _para_border_vals(p)
        para["shading"] = _para_shading(p)
        para["widow_control"] = _para_flag(p, "widowControl")
        para["top_line_punct"] = _para_flag(p, "topLinePunct")
        paragraphs.append(para)

    tables = []
    for i, t in enumerate(doc.tables):
        header = []
        try:
            header = [c.text.strip() for c in t.rows[0].cells][:8]
        except Exception:
            pass
        tables.append({
            "index": i, "rows": len(t.rows), "cols": len(t.columns),
            "align": _align_name(t.alignment), "header_row": header,
        })

    images = []
    try:
        for i, shp in enumerate(doc.inline_shapes):
            images.append({"index": i, "width_cm": _cm(shp.width),
                           "height_cm": _cm(shp.height), "wrap": "inline",
                           "position": None, "shadow": False})
    except Exception:
        pass
    try:
        # 浮动对象（w:drawing/wp:anchor）：图片计入 images，文本框文字计入 textbox_texts
        ns = {"wp": "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing",
              "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
              "pic": "http://schemas.openxmlformats.org/drawingml/2006/picture",
              "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
        textbox_texts = []
        for anc in doc.element.findall(".//w:drawing/wp:anchor", ns):
            wrap = None
            for wt in ("wp:wrapSquare", "wp:wrapTight", "wp:wrapTopAndBottom",
                       "wp:wrapNone", "wp:wrapThrough"):
                if anc.find(wt, ns) is not None:
                    wrap = wt.split(":")[1].replace("wrap", "").lower()

            def _pos(el):
                if el is None:
                    return None
                a = el.find("wp:align", ns)
                if a is not None:
                    return (el.get("relativeFrom"), a.text)
                o = el.find("wp:posOffset", ns)
                return (el.get("relativeFrom"), o.text if o is not None else None)

            ph = _pos(anc.find("wp:positionH", ns))
            pv = _pos(anc.find("wp:positionV", ns))
            pos = None
            if ph or pv:
                pos = "H:%s/%s V:%s/%s" % (
                    ph[0] if ph else "?", ph[1] if ph else "?",
                    pv[0] if pv else "?", pv[1] if pv else "?")
            shadow = anc.find(".//a:effectLst/a:outerShdw", ns) is not None
            # 文本框文字（题注等）：drawingml a:t 与 wordprocessing w:t 都收
            try:
                ttxt = "".join(t.text or "" for t in anc.findall(".//a:t", ns))
                ttxt += "".join(t.text or "" for t in anc.findall(".//w:t", ns))
                if ttxt.strip():
                    textbox_texts.append(ttxt.strip())
            except Exception:
                pass
            # 只有真图片（含 pic:pic）才计入 images，尺寸读 pic/spPr/xfrm/ext
            if anc.find(".//pic:pic", ns) is None:
                continue
            ext = anc.find(".//pic:spPr/a:xfrm/a:ext", ns)
            if ext is None:
                ext = anc.find(".//a:graphic/a:graphicData/pic:pic/pic:spPr/a:xfrm/a:ext", ns)
            w = h = None
            if ext is not None:
                w, h = _cm(ext.get("cx")), _cm(ext.get("cy"))
            images.append({"index": len(images), "width_cm": w, "height_cm": h,
                           "wrap": wrap, "position": pos, "shadow": shadow})
    except Exception:
        textbox_texts = []

    sections = []
    try:
        for sec in doc.sections:
            s = {
                "margin_top_cm": _cm(sec.top_margin),
                "margin_bottom_cm": _cm(sec.bottom_margin),
                "margin_left_cm": _cm(sec.left_margin),
                "margin_right_cm": _cm(sec.right_margin),
                "paper_w_cm": _cm(sec.page_width),
                "paper_h_cm": _cm(sec.page_height),
                "orientation": "landscape" if sec.orientation == 1 else "portrait",
                "gutter_cm": _cm(sec.gutter),
                "header": sec.header.paragraphs[0].text if sec.header.paragraphs else "",
                "footer": sec.footer.paragraphs[0].text if sec.footer.paragraphs else "",
            }
            # 页眉突出显示（黄色等）
            try:
                hl = None
                for hp in sec.header.paragraphs:
                    for r in hp.runs:
                        v = _hl_color(r)
                        if v:
                            hl = v
                            break
                s["header_highlight"] = hl
            except Exception:
                pass
            # 页码格式/起始值（与 python-docx 同一 sectPr 源，避免取到别的节）
            try:
                sp = sec._sectPr
                pgmar = sp.find("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}pgMar")
                if pgmar is not None:
                    gp = pgmar.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}gutterPos")
                    # 修正（2026-08-19）：无 gutterPos 属性时 Word 缺省即 left，
                    # 显式写 "left" 避免原语判定落空（gutter_pos 原语据此返回 "left"）
                    s["gutter_pos"] = gp if gp else "left"
                pnt = sp.find("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}pgNumType")
                if pnt is not None:
                    fmt = pnt.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}fmt")
                    if fmt:
                        s["page_number_fmt"] = fmt
                    st = pnt.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}start")
                    if st is not None:
                        try:
                            s["page_number_start"] = int(st)
                        except (TypeError, ValueError):
                            pass
                # 分栏：w:cols num/sep（多栏排版判定）
                cols = sp.find("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}cols")
                if cols is not None:
                    try:
                        s["cols_num"] = int(cols.get(
                            "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}num", 1))
                    except (TypeError, ValueError):
                        s["cols_num"] = 1
                    s["cols_sep"] = cols.get(
                        "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}sep") == "1"
                # 页面边框：w:pgBorders 存在且四边 val 非 none
                pb = sp.find("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}pgBorders")
                if pb is not None:
                    sides = {}
                    for side in ("top", "bottom", "left", "right"):
                        el = pb.find("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}%s" % side)
                        if el is not None:
                            val = el.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val")
                            sides[side] = {
                                "val": val,
                                "color": el.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}color"),
                                "sz": el.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}sz"),
                                "space": el.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}space"),
                            }
                    s["page_border"] = {
                        "offsetFrom": pb.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}offsetFrom"),
                        "frames": sides,
                        "valid": bool(sides) and any(
                            v.get("val") and v.get("val") != "none" for v in sides.values()),
                    }
            except Exception:
                pass
            sections.append(s)
    except Exception:
        pass
    header_text = sections[0].get("header", "") if sections else ""
    footer_text = sections[0].get("footer", "") if sections else ""

    # OOXML 直读（python-docx 不覆盖的部分）：先解出 document.xml 再解析，
    # 否则 _xml_root(data) 对 docx 整体字节（zip 容器）会失败返回 None，
    # 导致批注/公式/页码/水印/修订/目录等全部静默跳过。
    comments, formulas, footnotes, endnotes = [], 0, 0, 0
    page_breaks, track_changes, watermark, toc = 0, False, False, False
    doc_xml = _word_part(zf, "word/document.xml")
    root = _xml_root(doc_xml)
    if root is not None:
        comments_xml = _word_part(zf, "word/comments.xml")
        if comments_xml:
            croot = _xml_root(comments_xml)
            if croot is not None:
                for c in croot.findall(".//w:comment", _W_NS):
                    txt = "".join(t.text or "" for t in c.findall(".//w:t", _W_NS))
                    if txt:
                        comments.append(txt)
        formulas = _count_xpath(root, ".//m:oMath", _W_NS)
        footnotes = _count_xpath(root, ".//w:footnoteReference", _W_NS)
        endnotes = _count_xpath(root, ".//w:endnoteReference", _W_NS)
        page_breaks = _count_xpath(root, ".//w:br[@w:type='page']", _W_NS)
        track_changes = (_count_xpath(root, ".//w:ins", _W_NS) > 0 or
                         _count_xpath(root, ".//w:del", _W_NS) > 0)
        watermark = _count_xpath(root, ".//w:watermark", _W_NS) > 0
        toc = (_count_xpath(root, ".//w:fldSimple[@w:instr[contains(.,'TOC')]]", _W_NS) > 0 or
               any("TOC" in (el.get("w:instr") or "") for el in
                   root.findall(".//w:fldChar", _W_NS) if False) or
               _count_xpath(root, ".//w:instrText[contains(.,'TOC')]", _W_NS) > 0)

    # 页脚直读（sdt 包裹时 python-docx 取不到正文）：文本 + PAGE 域检测
    footer_page_field = False
    try:
        for nm in zf.namelist():
            import re
            if not re.match(r"word/footer\d+\.xml$", nm):
                continue
            fx = _xml_root(zf.read(nm))
            if fx is None:
                continue
            if (_count_xpath(fx, ".//w:instrText[contains(.,'PAGE')]", _W_NS) > 0 or
                    _count_xpath(fx, ".//w:fldSimple[@w:instr[contains(.,'PAGE')]]", _W_NS) > 0):
                footer_page_field = True
    except Exception:
        pass

    try:
        pg_num_fmt = pg_num_start = None
        sectPr = root.find(".//w:sectPr", _W_NS) if root is not None else None
        if sectPr is not None:
            pnt = sectPr.find("w:pgNumType", _W_NS)
            if pnt is not None:
                pg_num_fmt = pnt.get("w:fmt")
                pg_num_start = pnt.get("w:start")
        if sections:
            if pg_num_fmt and "page_number_fmt" not in sections[0]:
                sections[0]["page_number_fmt"] = pg_num_fmt
            if pg_num_start is not None and "page_number_start" not in sections[0]:
                try:
                    sections[0]["page_number_start"] = int(pg_num_start)
                except (TypeError, ValueError):
                    pass
    except Exception:
        pass

    return ParsedDoc(
        "word", paragraphs=paragraphs, tables=tables, images=images,
        comments=comments, formulas=formulas, footnotes=footnotes,
        endnotes=endnotes, header_text=header_text, footer_text=footer_text,
        sections=sections, page_breaks=page_breaks,
        track_changes=track_changes, watermark=watermark, toc=toc,
        meta=doc,
    )._with_footer_page(footer_page_field)._with_textboxes(textbox_texts)


# --------------------------------------------------------------------------
# xlsx
# --------------------------------------------------------------------------

def _parse_xlsx(zf, data):
    from openpyxl import load_workbook

    wb = load_workbook(BytesIO(data), data_only=False)
    sheet_count = len(wb.sheetnames)
    # 用第一个工作表做抽样
    ws = wb.active
    merged = 0
    formulas = 0
    fonts, aligns = [], []
    has_border = False
    fills = set()
    wrap_text = False
    number_formats = set()
    try:
        merged = len(ws.merged_cells.ranges)
    except Exception:
        pass
    try:
        for row in ws.iter_rows(max_row=min(ws.max_row or 0, 50),
                                max_col=min(ws.max_column or 0, 30)):
            for cell in row:
                if cell.value is not None:
                    v = str(cell.value)
                    if v.startswith("="):
                        formulas += 1
                    try:
                        if cell.font and cell.font.size:
                            fonts.append(cell.font.size)
                        if cell.alignment and cell.alignment.horizontal:
                            aligns.append(cell.alignment.horizontal)
                        if cell.number_format:
                            nf = str(cell.number_format)
                            if nf != "General":
                                number_formats.add(nf)
                        if cell.alignment and cell.alignment.wrap_text:
                            wrap_text = True
                    except Exception:
                        pass
                try:
                    if cell.border:
                        for side in (cell.border.top, cell.border.bottom,
                                     cell.border.left, cell.border.right):
                            if side and side.style:
                                has_border = True
                                break
                    if cell.fill and cell.fill.start_color and cell.fill.start_color.rgb:
                        rgb = cell.fill.start_color.rgb
                        if rgb and str(rgb) not in ("00000000", "FFFFFFFF"):
                            fills.add(str(rgb))
                except Exception:
                    pass
    except Exception:
        pass
    freeze = None
    try:
        if ws.freeze_panes:
            freeze = str(ws.freeze_panes)
    except Exception:
        pass
    protected = False
    try:
        protected = bool(ws.protection.sheet)
    except Exception:
        pass
    tables = []
    try:
        for i, t in enumerate(ws.tables.values()):
            tables.append({"index": i, "name": t.displayName,
                           "rows": 1, "cols": 1})
    except Exception:
        pass
    cond_format = 0
    cond_format_types = set()
    try:
        for cf in ws.conditional_formatting:
            cond_format += 1
            rule = getattr(cf, "rules", None) or []
            for rr in rule:
                t = getattr(rr, "type", None)
                if t:
                    cond_format_types.add(str(t))
    except Exception:
        pass
    auto_filter = False
    try:
        auto_filter = ws.auto_filter.ref is not None
    except Exception:
        pass
    data_validations = 0
    try:
        data_validations = len(ws.data_validations.dataValidation) if ws.data_validations else 0
    except Exception:
        pass
    hyperlinks = 0
    try:
        hyperlinks = len(ws._hyperlinks) if ws._hyperlinks else 0
    except Exception:
        pass
    print_area = False
    try:
        print_area = ws.print_area is not None
    except Exception:
        pass
    # 列宽/行高抽样
    col_widths, row_heights = [], []
    try:
        for col in list(ws.column_dimensions.values())[:20]:
            if col.width:
                col_widths.append(round(col.width, 1))
    except Exception:
        pass
    try:
        for row in list(ws.row_dimensions.values())[:20]:
            if row.height:
                row_heights.append(round(row.height, 1))
    except Exception:
        pass

    paragraphs = [{"index": 0, "text": "工作表: %s" % "、".join(wb.sheetnames[:10]),
                   "style": None}]
    return ParsedDoc(
        "xlsx", paragraphs=paragraphs, tables=tables,
        sheet_count=sheet_count, meta=wb,
    )._with_xlsx(merged, formulas, fonts, aligns, freeze, protected,
                 cond_format, auto_filter, col_widths, row_heights,
                 has_border, fills, wrap_text, number_formats,
                 data_validations, hyperlinks, print_area, cond_format_types)


def _parse_ppt(zf, data):
    from pptx import Presentation

    prs = Presentation(BytesIO(data))
    slides = []
    slide_count = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        texts = []
        fonts = []
        aligns = []
        has_bullet = False
        has_hyperlink = False
        images = 0
        shape_count = 0
        has_title = False
        try:
            for shape in slide.shapes:
                shape_count += 1
                if shape.shape_type == 13:  # PICTURE
                    images += 1
                # 标题占位符（title 或带 title 语义的 placeholder）
                try:
                    is_title = (getattr(shape, "is_title", False)
                                or (shape.is_placeholder
                                    and getattr(shape.placeholder_format, "type", None) is not None
                                    and "TITLE" in str(getattr(shape.placeholder_format, "type", "")).upper()))
                    if is_title and getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                        has_title = True
                except Exception:
                    pass
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text.strip())
                            if para.alignment is not None:
                                aligns.append(_align_name(para.alignment))
                            try:
                                if para.level and para.level > 0:
                                    has_bullet = True
                            except Exception:
                                pass
                        for run in para.runs:
                            try:
                                if run.font.size:
                                    fonts.append(run.font.size.pt)
                            except Exception:
                                pass
                            if run.hyperlink and run.hyperlink.address:
                                has_hyperlink = True
        except Exception:
            pass
        # 切换/动画：Timing 元素存在即视为有动画
        has_animation = False
        try:
            timing = slide._element.find(
                "{http://schemas.openxmlformats.org/presentationml/2006/main}timing")
            has_animation = timing is not None
        except Exception:
            pass
        has_transition = False
        try:
            has_transition = slide._element.find(
                "{http://schemas.openxmlformats.org/presentationml/2006/main}transition") is not None
        except Exception:
            pass
        notes = ""
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text[:200]
        except Exception:
            pass
        slides.append({
            "index": i, "texts": texts[:20], "fonts": fonts[:20],
            "aligns": aligns[:20], "bullet": has_bullet,
            "hyperlink": has_hyperlink, "images": images,
            "animation": has_animation, "transition": has_transition,
            "notes": notes, "has_title": has_title, "shape_count": shape_count,
        })
    paragraphs = [{"index": i, "text": s["texts"][0] if s["texts"] else "",
                   "style": None} for i, s in enumerate(slides)]
    # 文档级聚合：版式使用集合、比例（16:9 / 4:3）
    layouts = set()
    try:
        for slide in prs.slides:
            try:
                lname = (slide.slide_layout.name if slide.slide_layout else None)
                if lname:
                    layouts.add(str(lname))
            except Exception:
                pass
    except Exception:
        pass
    aspect = None
    try:
        sw, sh = prs.slide_width, prs.slide_height
        if sw and sh and sh:
            ratio = float(sw) / float(sh)
            aspect = "16:9" if ratio > 1.5 else ("4:3" if ratio < 1.5 else "other")
    except Exception:
        aspect = None
    return ParsedDoc(
        "pptx", paragraphs=paragraphs, slide_count=slide_count, meta=prs,
    )._with_slides(slides)._with_pptx_meta(layouts, aspect)


# 扩展字段（避免污染 ParsedDoc 主构造）
def _ext(self, **kw):
    for k, v in kw.items():
        setattr(self, k, v)
    return self


ParsedDoc._with_xlsx = lambda self, *a: _ext(
    self, xlsx_merged=a[0], xlsx_formulas=a[1], xlsx_fonts=a[2],
    xlsx_aligns=a[3], xlsx_freeze=a[4], xlsx_protected=a[5],
    xlsx_cond_format=a[6], xlsx_auto_filter=a[7],
    xlsx_col_widths=a[8], xlsx_row_heights=a[9],
    xlsx_border=a[10], xlsx_fills=a[11], xlsx_wrap_text=a[12],
    xlsx_number_formats=a[13], xlsx_data_validations=a[14],
    xlsx_hyperlinks=a[15], xlsx_print_area=a[16],
    xlsx_cond_format_types=a[17])
ParsedDoc._with_slides = lambda self, slides: _ext(self, slides=slides)
ParsedDoc._with_pptx_meta = lambda self, layouts, aspect: _ext(
    self, pptx_layouts=layouts, pptx_aspect=aspect)
ParsedDoc._with_footer_page = lambda self, v: _ext(self, footer_page_field=v)
ParsedDoc._with_textboxes = lambda self, v: _ext(self, textbox_texts=v)


def parse_bytes(data, file_type):
    if file_type not in ("word", "xlsx", "pptx"):
        raise ParseError("不支持的类型: %s" % file_type)
    zf = _check_zip(data, file_type)
    try:
        if file_type == "word":
            return _parse_word(zf, data)
        if file_type == "xlsx":
            return _parse_xlsx(zf, data)
        return _parse_ppt(zf, data)
    except ParseError:
        raise
    except Exception as e:
        raise ParseError("无法解析 %s 文档: %s" % (file_type, e)) from e


def parse_file(path):
    p = Path(path)
    ext = p.suffix.lower()
    ft = {".docx": "word", ".xlsx": "xlsx", ".pptx": "pptx"}.get(ext)
    if ft is None:
        raise ParseError("不支持的文件类型: %s（支持 .docx/.xlsx/.pptx）" % ext)
    return parse_bytes(p.read_bytes(), ft)


if __name__ == "__main__":
    import json
    import sys

    for f in sys.argv[1:]:
        try:
            d = parse_file(f)
            info = {k: getattr(d, k) for k in
                    ("file_type", "text", "tables", "images", "comments",
                     "formulas", "sections", "page_breaks", "sheet_count",
                     "slide_count")}
            print(json.dumps(info, ensure_ascii=False, indent=1)[:4000])
        except ParseError as e:
            print("PARSE_ERROR:", e)
